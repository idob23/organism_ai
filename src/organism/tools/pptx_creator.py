
from dataclasses import dataclass
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsmap
from typing import Any
from .base import BaseTool, ToolResult, OUTPUTS_DIR
import re


@dataclass
class Theme:
    bg: RGBColor
    title: RGBColor
    text: RGBColor
    subtitle: RGBColor
    accent: RGBColor


_LIGHT = Theme(
    bg=RGBColor(0xFF, 0xFF, 0xFF),
    title=RGBColor(0x1E, 0x3A, 0x5F),
    text=RGBColor(0x33, 0x33, 0x33),
    subtitle=RGBColor(0x66, 0x66, 0x66),
    accent=RGBColor(0x2E, 0x5C, 0x8A),
)

_DARK = Theme(
    bg=RGBColor(0x1A, 0x1A, 0x2E),
    title=RGBColor(0xFF, 0xFF, 0xFF),
    text=RGBColor(0xE0, 0xE0, 0xE0),
    subtitle=RGBColor(0xB0, 0xB8, 0xC8),
    accent=RGBColor(0xE9, 0x4F, 0x37),
)

_THEMES = {"light": _LIGHT, "dark": _DARK}


def _set_bg(slide, color: RGBColor):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _add_title_box(slide, text, left, top, width, height,
                   size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color


def _set_bullet(paragraph, char="\u2022"):
    """Set a bullet character on a paragraph via XML."""
    pPr = paragraph._p.get_or_add_pPr()
    # Remove buNone if present
    for child in list(pPr):
        if child.tag.endswith("}buNone"):
            pPr.remove(child)
    # Add bullet char
    bu = parse_xml(f'<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="{char}"/>')
    pPr.append(bu)


def _add_content_box(slide, text, left, top, width, height,
                     size=18, color=None):
    """Add a text box with parsed bullets and numbered lists."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True

    lines = text.split("\n")
    first = True
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.space_before = Pt(4)
        p.space_after = Pt(4)

        # Bullet line
        bullet_match = re.match(r'^[\u2022\-\*\u2014]\s*(.+)$', stripped)
        numbered_match = re.match(r'^(\d+)[.)]\s*(.+)$', stripped)

        if bullet_match:
            r = p.add_run()
            r.text = bullet_match.group(1)
            r.font.size = Pt(size)
            if color:
                r.font.color.rgb = color
            p.level = 0
            _set_bullet(p)
        elif numbered_match:
            r = p.add_run()
            r.text = f"{numbered_match.group(1)}. {numbered_match.group(2)}"
            r.font.size = Pt(size)
            if color:
                r.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = stripped
            r.font.size = Pt(size)
            if color:
                r.font.color.rgb = color


def _add_line(slide, left, top, width, color: RGBColor, height=0.02):
    s = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


class PptxCreatorTool(BaseTool):

    @property
    def name(self) -> str:
        return "pptx_creator"

    @property
    def description(self) -> str:
        return (
            "Create a PowerPoint presentation (.pptx file). "
            "Provide slides with title and content. "
            "Content can be brief \u2014 it will be expanded automatically. "
            "Supports light (default) and dark themes. "
            "Use bullet points (\u2022, -, *) in content for proper PowerPoint bullets."
        )

    @property
    def input_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "filename": {"type": "string"},
                "topic": {"type": "string", "description": "Presentation topic for context"},
                "theme": {
                    "type": "string",
                    "enum": ["light", "dark"],
                    "default": "light",
                    "description": "Presentation theme (light for projector, dark for screen)",
                },
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "content": {"type": "string"},
                            "notes": {
                                "type": "string",
                                "description": "Speaker notes for this slide",
                            },
                        }
                    }
                }
            },
            "required": ["filename", "slides"],
        }

    async def _expand_content(self, topic: str, slide_title: str, brief: str) -> str:
        if len(brief) > 300:
            return brief
        from src.organism.llm.claude import ClaudeProvider
        from src.organism.llm.base import Message
        llm = ClaudeProvider()
        response = await llm.complete(
            messages=[Message(role="user", content=(
                f"\u0422\u0435\u043c\u0430 \u043f\u0440\u0435\u0437\u0435\u043d\u0442\u0430\u0446\u0438\u0438: {topic}\n"
                f"\u0417\u0430\u0433\u043e\u043b\u043e\u0432\u043e\u043a \u0441\u043b\u0430\u0439\u0434\u0430: {slide_title}\n"
                f"\u041a\u0440\u0430\u0442\u043a\u043e\u0435 \u043e\u043f\u0438\u0441\u0430\u043d\u0438\u0435: {brief}\n\n"
                f"\u041d\u0430\u043f\u0438\u0448\u0438 \u0441\u043e\u0434\u0435\u0440\u0436\u0438\u043c\u043e\u0435 \u0441\u043b\u0430\u0439\u0434\u0430: "
                f"4-6 \u0442\u0435\u0437\u0438\u0441\u043e\u0432 \u0447\u0435\u0440\u0435\u0437 bullet points "
                f"(\u0441\u0438\u043c\u0432\u043e\u043b \u2022). "
                f"\u041a\u0430\u0436\u0434\u044b\u0439 \u0442\u0435\u0437\u0438\u0441 \u2014 "
                f"\u043f\u043e\u043b\u043d\u043e\u0435 \u043f\u0440\u0435\u0434\u043b\u043e\u0436\u0435\u043d\u0438\u0435, "
                f"\u043d\u0435 \u043e\u0431\u0440\u044b\u0432\u0430\u0439 \u043c\u044b\u0441\u043b\u044c. "
                f"\u042f\u0437\u044b\u043a: \u0440\u0443\u0441\u0441\u043a\u0438\u0439. "
                f"\u041c\u0430\u043a\u0441\u0438\u043c\u0443\u043c 800 \u0441\u0438\u043c\u0432\u043e\u043b\u043e\u0432."
            ))],
            system="\u041f\u0438\u0448\u0438 \u0442\u043e\u043b\u044c\u043a\u043e \u0442\u0435\u043a\u0441\u0442 \u0441\u043b\u0430\u0439\u0434\u0430. "
                   "\u041d\u0438\u043a\u0430\u043a\u0438\u0445 \u0432\u0441\u0442\u0443\u043f\u043b\u0435\u043d\u0438\u0439. "
                   "\u041d\u0430\u0447\u0438\u043d\u0430\u0439 \u0441\u0440\u0430\u0437\u0443 \u0441 \u0441\u043e\u0434\u0435\u0440\u0436\u0430\u043d\u0438\u044f.",
            model_tier="balanced",
            max_tokens=1000,
        )
        return response.content.strip()

    async def execute(self, input: dict[str, Any]) -> ToolResult:
        filename: str = input["filename"]
        slides_data: list[dict] = input["slides"]
        topic: str = input.get("topic", slides_data[0].get("title", "") if slides_data else "")
        theme_name: str = input.get("theme", "light")
        theme = _THEMES.get(theme_name, _LIGHT)

        if not filename.endswith(".pptx"):
            filename += ".pptx"

        expanded = []
        for slide in slides_data:
            content = slide.get("content", "")
            title = slide.get("title", "")
            notes = slide.get("notes", "")
            if len(content) < 300 and title and title != slides_data[0].get("title"):
                content = await self._expand_content(topic, title, content or title)
            expanded.append({"title": title, "content": content, "notes": notes})

        try:
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            for i, s in enumerate(expanded):
                layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(layout)
                _set_bg(slide, theme.bg)

                if i == 0:
                    # Title slide
                    _add_title_box(slide, s["title"], 0.8, 2.3, 11.5, 1.5,
                                   size=36, bold=True, color=theme.title,
                                   align=PP_ALIGN.CENTER)
                    if topic and topic != s["title"]:
                        _add_title_box(slide, topic, 0.8, 4.0, 11.5, 1.0,
                                       size=20, color=theme.subtitle,
                                       align=PP_ALIGN.CENTER)
                    elif s["content"]:
                        _add_title_box(slide, s["content"], 0.8, 4.0, 11.5, 1.5,
                                       size=20, color=theme.subtitle,
                                       align=PP_ALIGN.CENTER)
                    _add_line(slide, 2.0, 5.5, 9.33, theme.accent, height=0.03)
                else:
                    # Content slide
                    _add_line(slide, 0, 0, 13.33, theme.accent, height=0.05)
                    _add_title_box(slide, s["title"], 0.5, 0.2, 12.3, 0.9,
                                   size=28, bold=True, color=theme.title)
                    _add_line(slide, 0.5, 1.15, 12.3, theme.subtitle, height=0.015)
                    if s["content"]:
                        _add_content_box(slide, s["content"], 0.5, 1.35, 12.3, 5.5,
                                         size=18, color=theme.text)
                    # Slide number
                    _add_title_box(slide, str(i + 1), 12.4, 6.9, 0.7, 0.4,
                                   size=12, color=theme.subtitle, align=PP_ALIGN.RIGHT)

                # Speaker notes
                if s.get("notes"):
                    notes_slide = slide.notes_slide
                    notes_tf = notes_slide.notes_text_frame
                    notes_tf.text = s["notes"]

            filepath = OUTPUTS_DIR / Path(filename).name
            prs.save(str(filepath))
            _fname = Path(filename).name
            return ToolResult(output=f"Saved files: {_fname}", exit_code=0,
                              created_files=[_fname])

        except Exception as e:
            return ToolResult(output="", error=str(e), exit_code=1)
